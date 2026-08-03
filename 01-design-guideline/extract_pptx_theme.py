"""
Extrahiert Design-Tokens aus einem PowerPoint-Folienmaster (.potx oder .pptx):
Theme-Farben, Schriftfamilien, Layout-Namen, Platzhalter-Positionen,
Textstile je Gliederungsebene und Hintergrundfuellung je Layout.

Benoetigt: pip install python-pptx

Aufruf:
    python extract_pptx_theme.py "Folienmaster.potx"

Output:
    pptx_design_tokens.json  (im selben Ordner)

Farbwerte erscheinen entweder als "#RRGGBB" (feste Farbe) oder als
"scheme:<name>" (Verweis auf eine Theme-Farbe, siehe theme_colors).
Fontangaben "+mj-lt" bzw. "+mn-lt" verweisen auf heading_font bzw. body_font
aus theme_fonts.
"""

import json
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

# DrawingML-Fuellarten, wie sie unter <p:bgPr> vorkommen koennen
FILL_TAGS = {
    "solidFill": "solid",
    "gradFill": "gradient",
    "blipFill": "picture",
    "pattFill": "pattern",
    "grpFill": "group",
    "noFill": "none",
}


def emu_to_cm(value):
    if value is None:
        return None
    return round(Emu(value).cm, 2)


def _color_of(element):
    """Erste Farbangabe unterhalb von element: "#RRGGBB" oder "scheme:<name>"."""
    if element is None:
        return None
    srgb = element.find(f".//{qn('a:srgbClr')}")
    if srgb is not None:
        return f"#{srgb.get('val', '').upper()}"
    scheme = element.find(f".//{qn('a:schemeClr')}")
    if scheme is not None:
        return f"scheme:{scheme.get('val')}"
    return None


def extract_theme_colors(prs):
    """Liest die 12 Theme-Farben (Accent1-6, Dark1/2, Light1/2, Hyperlink,
    FollowedHyperlink)."""
    colors = {}
    try:
        # python-pptx legt Theme-Farben nicht direkt offen; über den Theme-Part gehen
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        xml = theme_part.blob.decode("utf-8", errors="ignore")

        import re

        # sucht srgbClr Werte innerhalb der clrScheme
        scheme_match = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.DOTALL)
        if scheme_match:
            scheme_xml = scheme_match.group(0)
            entries = re.findall(
                r'<a:(\w+)>\s*<a:srgbClr val="([0-9A-Fa-f]{6})"', scheme_xml
            )
            for name, hexval in entries:
                colors[name] = f"#{hexval.upper()}"
            # sysClr (z.B. für dk1/lt1 die windowText/window referenzieren)
            sys_entries = re.findall(
                r'<a:(\w+)>\s*<a:sysClr val="(\w+)" lastClr="([0-9A-Fa-f]{6})"',
                scheme_xml,
            )
            for name, sysname, hexval in sys_entries:
                colors[name] = f"#{hexval.upper()} (sysClr: {sysname})"
    except Exception as e:
        colors["_error"] = str(e)
    return colors


def extract_theme_fonts(prs):
    fonts = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        xml = theme_part.blob.decode("utf-8", errors="ignore")

        import re

        major = re.search(
            r'<a:majorFont>.*?<a:latin typeface="([^"]+)"', xml, re.DOTALL
        )
        minor = re.search(
            r'<a:minorFont>.*?<a:latin typeface="([^"]+)"', xml, re.DOTALL
        )
        fonts["heading_font"] = major.group(1) if major else None
        fonts["body_font"] = minor.group(1) if minor else None
    except Exception as e:
        fonts["_error"] = str(e)
    return fonts


def _style_levels(style_element):
    """Liest lvl1pPr..lvl9pPr eines txStyles-Blocks (titleStyle, bodyStyle, ...)."""
    levels = []
    if style_element is None:
        return levels
    for level in range(1, 10):
        lvl = style_element.find(qn(f"a:lvl{level}pPr"))
        if lvl is None:
            continue
        def_rpr = lvl.find(qn("a:defRPr"))
        size = def_rpr.get("sz") if def_rpr is not None else None
        latin = def_rpr.find(qn("a:latin")) if def_rpr is not None else None
        levels.append(
            {
                "level": level,
                "font_size_pt": round(int(size) / 100, 1) if size else None,
                "color": (
                    _color_of(def_rpr.find(qn("a:solidFill")))
                    if def_rpr is not None
                    else None
                ),
                "align": lvl.get("algn"),
                "font": latin.get("typeface") if latin is not None else None,
            }
        )
    return levels


def extract_text_styles(master):
    """Schriftgroesse, Farbe, Ausrichtung und Font je Gliederungsebene aus den
    txStyles des Folienmasters."""
    styles = {"title": [], "body": []}
    tx_styles = master.element.find(qn("p:txStyles"))
    if tx_styles is None:
        return styles
    styles["title"] = _style_levels(tx_styles.find(qn("p:titleStyle")))
    styles["body"] = _style_levels(tx_styles.find(qn("p:bodyStyle")))
    return styles


def extract_background(slide_element):
    """Hintergrundfuellung eines Layouts oder Masters.

    fill = "inherited" heisst: kein eigenes <p:bg>, der Hintergrund kommt vom
    Master (bei Layouts) bzw. vom Theme (beim Master).
    """
    empty = {"fill": "inherited", "color": None}
    c_sld = slide_element.find(qn("p:cSld"))
    if c_sld is None:
        return empty
    bg = c_sld.find(qn("p:bg"))
    if bg is None:
        return empty

    bg_pr = bg.find(qn("p:bgPr"))
    if bg_pr is not None:
        for tag, fill_name in FILL_TAGS.items():
            fill = bg_pr.find(qn(f"a:{tag}"))
            if fill is not None:
                return {"fill": fill_name, "color": _color_of(fill)}
        return {"fill": None, "color": None}

    # <p:bgRef> verweist auf einen Fuellstil des Themes plus eine Farbe
    bg_ref = bg.find(qn("p:bgRef"))
    if bg_ref is not None:
        return {"fill": "theme_ref", "color": _color_of(bg_ref)}
    return {"fill": None, "color": None}


def extract_layouts(prs):
    layouts = []
    for master_idx, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                placeholders.append(
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "type_name": ph_type.name if ph_type is not None else None,
                        "type_id": ph_type.value if ph_type is not None else None,
                        "name": ph.name,
                        "left_cm": emu_to_cm(ph.left),
                        "top_cm": emu_to_cm(ph.top),
                        "width_cm": emu_to_cm(ph.width),
                        "height_cm": emu_to_cm(ph.height),
                    }
                )
            layouts.append(
                {
                    "master_index": master_idx,
                    "layout_name": layout.name,
                    "background": extract_background(layout.element),
                    "placeholders": placeholders,
                }
            )
    return layouts


def extract(pptx_path: str):
    prs = Presentation(pptx_path)

    slide_width_cm = emu_to_cm(prs.slide_width)
    slide_height_cm = emu_to_cm(prs.slide_height)

    # ponytail: Textstile und Master-Hintergrund vom ersten Master, wie schon
    # theme_colors und theme_fonts. Erst wenn eine Datei mehrere Master mit
    # unterschiedlichen txStyles hat, lohnt eine Liste je master_index.
    master = prs.slide_masters[0]

    result = {
        "source_file": pptx_path,
        "slide_width_cm": slide_width_cm,
        "slide_height_cm": slide_height_cm,
        "aspect_ratio": (
            round(prs.slide_width / prs.slide_height, 3) if prs.slide_height else None
        ),
        "theme_colors": extract_theme_colors(prs),
        "theme_fonts": extract_theme_fonts(prs),
        "text_styles": extract_text_styles(master),
        "master_background": extract_background(master.element),
        "layouts": extract_layouts(prs),
    }
    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_theme.py <pptx-oder-potx-datei>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    data = extract(pptx_path)

    out_path = "pptx_design_tokens.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    print(f"Fertig. Ergebnis geschrieben nach: {out_path}")
    print(f"Folienformat: {data['slide_width_cm']} x {data['slide_height_cm']} cm")
    print(f"Theme-Farben: {data['theme_colors']}")
    print(f"Theme-Fonts: {data['theme_fonts']}")
    print(
        "Textstile: "
        f"{len(data['text_styles']['title'])} Titel-Ebenen, "
        f"{len(data['text_styles']['body'])} Body-Ebenen"
    )
    print(f"Anzahl Layouts gefunden: {len(data['layouts'])}")
