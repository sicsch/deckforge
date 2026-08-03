"""Tests fuer die Folienmaster-Extraktion (Issue #77).

Die Test-Praesentation wird zur Laufzeit erzeugt und enthaelt bewusst keinerlei
Design-Werte eines konkreten Unternehmens (CLAUDE.md Regeln 1 und 4).
"""

import importlib.util
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

MODULE_PATH = (
    Path(__file__).resolve().parents[1]
    / "01-design-guideline"
    / "extract_pptx_theme.py"
)
# Der Ordnername beginnt mit einer Ziffer, also nicht importierbar -> ueber den Pfad
_spec = importlib.util.spec_from_file_location("extract_pptx_theme", MODULE_PATH)
extract_pptx_theme = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(extract_pptx_theme)

BG_COLOR = "1A73E8"  # generischer Platzhalter, kein reales Corporate Design


@pytest.fixture(scope="module")
def tokens(tmp_path_factory):
    """Synthetische .pptx mit einem eingefaerbten Layout, dann extrahiert."""
    prs = Presentation()
    layout = prs.slide_masters[0].slide_layouts[0]
    c_sld = layout.element.find(qn("p:cSld"))
    c_sld.insert(
        0,
        parse_xml(
            f"<p:bg {nsdecls('p', 'a')}><p:bgPr>"
            f'<a:solidFill><a:srgbClr val="{BG_COLOR}"/></a:solidFill>'
            f"<a:effectLst/></p:bgPr></p:bg>"
        ),
    )

    path = tmp_path_factory.mktemp("pptx") / "synthetic.pptx"
    prs.save(path)
    return extract_pptx_theme.extract(str(path))


def test_text_styles_have_size_and_color_per_level(tokens):
    styles = tokens["text_styles"]

    assert styles["title"] and styles["body"]
    for levels in (styles["title"], styles["body"]):
        for level in levels:
            assert level["level"] >= 1
            assert isinstance(level["font_size_pt"], float)
            assert level["color"]
    # Gliederungsebenen sind luecken- und duplikatfrei durchnummeriert
    assert [lvl["level"] for lvl in styles["body"]] == list(
        range(1, len(styles["body"]) + 1)
    )


def test_background_is_read_per_layout(tokens):
    first, *rest = tokens["layouts"]

    assert first["background"] == {"fill": "solid", "color": f"#{BG_COLOR}"}
    # Layouts ohne eigenes <p:bg> erben vom Master
    assert all(layout["background"]["fill"] == "inherited" for layout in rest)
    assert "fill" in tokens["master_background"]


def test_placeholder_types_are_machine_readable(tokens):
    placeholders = [ph for lo in tokens["layouts"] for ph in lo["placeholders"]]

    assert placeholders
    for ph in placeholders:
        assert isinstance(ph["type_id"], int)
        assert ph["type_name"].isupper()
        # der alte String bleibt die Kombination aus beiden
        assert ph["type"] == f"{ph['type_name']} ({ph['type_id']})"

    assert {"TITLE", "BODY", "PICTURE"} <= {ph["type_name"] for ph in placeholders}


def test_existing_fields_stay_unchanged(tokens):
    assert {
        "source_file",
        "slide_width_cm",
        "slide_height_cm",
        "aspect_ratio",
        "theme_colors",
        "theme_fonts",
        "layouts",
    } <= set(tokens)

    layout = tokens["layouts"][0]
    assert {"master_index", "layout_name", "placeholders"} <= set(layout)
    assert {
        "idx",
        "type",
        "name",
        "left_cm",
        "top_cm",
        "width_cm",
        "height_cm",
    } <= set(layout["placeholders"][0])
