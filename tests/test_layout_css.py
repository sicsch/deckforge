"""Tests fuer die deterministische Layout-CSS-Erzeugung (Issue #78).

Die Token-JSON wird zur Laufzeit aus einer synthetischen Praesentation erzeugt
und enthaelt keinerlei Design-Werte eines konkreten Unternehmens
(CLAUDE.md Regeln 1 und 4).
"""

import importlib.util
import re
from pathlib import Path

import layout_css
import pytest
from pptx import Presentation

MODULE_PATH = (
    Path(__file__).resolve().parents[1]
    / "01-design-guideline"
    / "extract_pptx_theme.py"
)
_spec = importlib.util.spec_from_file_location("extract_pptx_theme", MODULE_PATH)
extract_pptx_theme = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(extract_pptx_theme)


@pytest.fixture(scope="module")
def tokens(tmp_path_factory):
    path = tmp_path_factory.mktemp("pptx") / "synthetic.pptx"
    Presentation().save(path)
    return extract_pptx_theme.extract(str(path))


def test_cm_to_percent_self_check():
    """Selbstpruefung der cm-nach-Prozent-Umrechnung."""
    assert layout_css.cm_to_percent(0, 33.87) == 0.0
    assert layout_css.cm_to_percent(33.87, 33.87) == 100.0
    assert layout_css.cm_to_percent(16.935, 33.87) == 50.0
    assert layout_css.cm_to_percent(1.27, 25.4) == 5.0
    # Fehlende Werte erzeugen keine Regel statt einer falschen
    assert layout_css.cm_to_percent(None, 33.87) is None
    assert layout_css.cm_to_percent(1.0, None) is None
    assert layout_css.cm_to_percent(1.0, 0) is None


def test_placeholder_geometry_matches_master(tokens):
    """Abweichung der Platzhalterboxen unter 1 % der Foliengroesse."""
    width_cm = tokens["slide_width_cm"]
    height_cm = tokens["slide_height_cm"]
    by_name = {layout["layout_name"]: layout for layout in tokens["layouts"]}

    for entry in layout_css.layouts(tokens):
        source = by_name[entry["name"]]
        for slot, ph in zip(entry["slots"], source["placeholders"], strict=True):
            for key, cm_key, total in (
                ("left_pct", "left_cm", width_cm),
                ("top_pct", "top_cm", height_cm),
                ("width_pct", "width_cm", width_cm),
                ("height_pct", "height_cm", height_cm),
            ):
                back_cm = slot[key] / 100 * total
                assert abs(back_cm - ph[cm_key]) < total / 100


def test_style_block_is_complete(tokens):
    css = layout_css.build_css(tokens)

    assert css.startswith("<style>") and css.endswith("</style>")
    assert ":root {" in css
    assert "--slide-width: 1920px;" in css
    assert "--font-body:" in css and "--color-accent1:" in css
    assert "--title-size-1:" in css and "--body-size-1:" in css
    assert ".slide {" in css and ".ph {" in css
    # keine Restwerte in cm, Geometrie ausschliesslich in Prozent
    assert "cm;" not in css
    assert re.search(r"\.layout-[\w-]+ \.slot-[\w-]+ \{\n  left: [\d.]+%;", css)


def test_layout_selection_filters_css(tokens):
    everything = layout_css.layouts(tokens)
    first = everything[0]["slug"]

    selected = layout_css.layouts(tokens, [first])
    assert [entry["slug"] for entry in selected] == [first]

    css = layout_css.build_css(tokens, [first])
    assert f".layout-{first} " in css
    for entry in everything[1:]:
        assert f".{entry['css_class']} " not in css


def test_components_only_apply_inside_placeholders(tokens):
    css = layout_css.build_css(tokens)

    for name in layout_css.COMPONENT_CLASSES:
        for match in re.finditer(rf"^\.[^\n{{]*\.{re.escape(name)}\b", css, re.M):
            assert match.group().startswith(".ph "), match.group()


def test_catalog_lists_layouts_and_classes(tokens):
    markdown = layout_css.catalog_markdown(tokens)
    css = layout_css.build_css(tokens)

    for entry in layout_css.layouts(tokens):
        assert entry["name"] in markdown
        assert entry["css_class"] in markdown
    for name in layout_css.COMPONENT_CLASSES:
        assert f"`{name}`" in markdown
        assert f".{name}" in css


def test_master_css_is_stripped_and_restored_verbatim(tokens):
    """Iterationen sehen den generierten Block nicht und veraendern ihn nicht."""
    block = layout_css.build_css(tokens)
    deck = (
        f"<html><head>{block}<style>.x{{color:red}}</style></head>"
        "<body>D</body></html>"
    )

    stripped, cut = layout_css.split_master_css(deck)

    assert cut == block
    assert layout_css.MASTER_CSS_MARKER not in stripped
    assert layout_css.MASTER_CSS_PLACEHOLDER in stripped
    assert ".x{color:red}" in stripped  # eigenes CSS des Modells bleibt drin
    assert layout_css.restore_master_css(stripped, cut) == deck


def test_master_css_returns_even_if_the_model_drops_the_placeholder(tokens):
    block = layout_css.build_css(tokens)
    restored = layout_css.restore_master_css("<html><head></head></html>", block)
    assert block in restored

    # Modell hat den Block selbst nachgebaut: Original gewinnt, nur einmal drin
    rebuilt = (
        f"<html><head><style>\n{layout_css.MASTER_CSS_MARKER}\n"
        ".slide{}</style></head></html>"
    )
    restored = layout_css.restore_master_css(rebuilt, block)
    assert block in restored
    assert restored.count(layout_css.MASTER_CSS_MARKER) == 1


def test_deck_without_master_css_survives_the_roundtrip():
    deck = "<html><body>Kein Folienmaster</body></html>"
    stripped, cut = layout_css.split_master_css(deck)
    assert (stripped, cut) == (deck, "")
    assert layout_css.restore_master_css(deck, cut) == deck


def test_empty_tokens_stay_harmless():
    assert layout_css.layouts({}) == []
    assert "Keine Folienlayouts" in layout_css.catalog_markdown({})
    assert layout_css.build_css({}).startswith("<style>")
